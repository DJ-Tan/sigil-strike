"""
build_deck.py
─────────────
Generates the Sigil Strike team-briefing PowerPoint deck.

Run:
    python scripts/build_deck.py
Output:
    sigil_strike_team_briefing.pptx  (at the repo root)
"""

from __future__ import annotations

import pathlib

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


# ── Theme ───────────────────────────────────────────────────────────────────────
BG       = RGBColor(0x0F, 0x14, 0x2A)   # deep navy
PANEL    = RGBColor(0x1A, 0x22, 0x3D)
ACCENT   = RGBColor(0xFF, 0xB8, 0x4D)   # warm orange
ACCENT2  = RGBColor(0x6E, 0xC1, 0xE4)   # cool cyan
GOOD     = RGBColor(0x7E, 0xD3, 0x21)
BAD      = RGBColor(0xFF, 0x5C, 0x5C)
TEXT     = RGBColor(0xF2, 0xF2, 0xF7)
MUTED    = RGBColor(0xA8, 0xB0, 0xC4)

SLIDE_W  = Inches(13.333)
SLIDE_H  = Inches(7.5)


def _bg(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _set_text(tf, text, *, size=18, bold=False, color=TEXT, align=PP_ALIGN.LEFT,
              font="Calibri"):
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    p.text = ""
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def _add_text(slide, left, top, width, height, text, **kwargs):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    _set_text(tf, text, **kwargs)
    return box


def _add_bullets(slide, left, top, width, height, items, *, size=18,
                 color=TEXT, accent_color=ACCENT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(8)
        run = p.add_run()
        run.text = "▸  "
        run.font.name = "Calibri"
        run.font.size = Pt(size)
        run.font.bold = True
        run.font.color.rgb = accent_color
        run2 = p.add_run()
        run2.text = item
        run2.font.name = "Calibri"
        run2.font.size = Pt(size)
        run2.font.color.rgb = color
    return box


def _panel(slide, left, top, width, height, color=PANEL):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def _accent_bar(slide, left, top, width, height, color=ACCENT):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def _slide_header(slide, title, subtitle=None):
    _accent_bar(slide, Inches(0.5), Inches(0.45), Inches(0.12), Inches(0.7))
    _add_text(slide, Inches(0.75), Inches(0.35), Inches(11.5), Inches(0.8),
              title, size=32, bold=True, color=TEXT)
    if subtitle:
        _add_text(slide, Inches(0.75), Inches(0.95), Inches(11.5), Inches(0.5),
                  subtitle, size=14, color=MUTED)


def _footer(slide, num, total):
    _add_text(slide, Inches(11.0), Inches(7.05), Inches(2.0), Inches(0.3),
              f"SIGIL STRIKE   {num} / {total}", size=10, color=MUTED,
              align=PP_ALIGN.RIGHT)


# ── Slide builders ──────────────────────────────────────────────────────────────

def slide_title(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s)
    _accent_bar(s, Inches(0), Inches(3.0), Inches(13.333), Inches(0.05))
    _accent_bar(s, Inches(0), Inches(4.7),  Inches(13.333), Inches(0.05),
                color=ACCENT2)

    _add_text(s, Inches(0.5), Inches(3.1), Inches(12.3), Inches(1.4),
              "SIGIL STRIKE", size=80, bold=True, color=ACCENT,
              align=PP_ALIGN.CENTER)
    _add_text(s, Inches(0.5), Inches(4.15), Inches(12.3), Inches(0.6),
              "Two-Player Hand-Sign Fighter", size=28, color=TEXT,
              align=PP_ALIGN.CENTER)
    _add_text(s, Inches(0.5), Inches(4.85), Inches(12.3), Inches(0.5),
              "DSO World of Science 2026 — AI Computer Vision Workshop",
              size=18, color=MUTED, align=PP_ALIGN.CENTER)
    _add_text(s, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.4),
              "Team Briefing & Preparation Guide", size=14, color=ACCENT2,
              align=PP_ALIGN.CENTER)
    return s


def slide_what_is_it(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s)
    _slide_header(s, "What is Sigil Strike?",
                  "A real-time 1v1 fighter controlled by hand gestures")

    _panel(s, Inches(0.6), Inches(1.7), Inches(6.0), Inches(5.2))
    _add_text(s, Inches(0.9), Inches(1.9), Inches(5.6), Inches(0.5),
              "The Concept", size=22, bold=True, color=ACCENT)
    _add_bullets(s, Inches(0.9), Inches(2.5), Inches(5.6), Inches(4.5), [
        "Two players face the webcam and throw hand signs at the camera",
        "An AI model (trained by your team!) classifies each gesture",
        "Three gestures in a row = one move (attack / defend / heal)",
        "Every 5 seconds the game resolves one move per player",
        "First player to drop the opponent's HP to zero wins the round",
        "Keyboard input always works as a fallback if the model misfires",
    ], size=16)

    _panel(s, Inches(6.9), Inches(1.7), Inches(5.9), Inches(5.2))
    _add_text(s, Inches(7.2), Inches(1.9), Inches(5.5), Inches(0.5),
              "Why It Matters", size=22, bold=True, color=ACCENT2)
    _add_bullets(s, Inches(7.2), Inches(2.5), Inches(5.5), Inches(4.5), [
        "Each team trains its OWN model — you own the prediction quality",
        "Two pipelines to choose from: lightweight Landmark or deeper CNN",
        "Real end-to-end ML: capture → train → evaluate → deploy → compete",
        "Tournament uses your model live — accuracy and latency both matter",
        "Your weights ship next to the game; no rebuild needed",
    ], size=16, accent_color=ACCENT2)


def slide_tournament_format(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s)
    _slide_header(s, "Tournament Format",
                  "Six teams · two groups · knockouts · final")

    _panel(s, Inches(0.6), Inches(1.7), Inches(12.2), Inches(5.2))

    _add_text(s, Inches(0.9), Inches(1.95), Inches(11.6), Inches(0.5),
              "Group Stage  →  Semi-Finals  →  3rd-Place & Final",
              size=24, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

    _add_bullets(s, Inches(1.2), Inches(2.8), Inches(11.0), Inches(4.0), [
        "Six teams split into two groups of three (Group A & Group B)",
        "Round-robin within each group — every team plays every other in its group",
        "Top two from each group advance to the semi-finals (cross-bracket)",
        "5th-place playoff for the two 3rd-placed teams",
        "Semi-final losers meet in the 3rd-place playoff",
        "Semi-final winners meet in the Grand Final",
        "Match length: HP-based — first player to 0 HP loses; deathmatch ticks if time runs long",
    ], size=17)


def slide_gameplay_loop(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s)
    _slide_header(s, "Core Gameplay Loop",
                  "Action → Buffer → Combo → Resolve")

    # Four-step horizontal flow
    steps = [
        ("1. SIGN",   "Throw a hand sign at the camera (or press a key).",        ACCENT),
        ("2. BUFFER", "Your gesture joins a rolling buffer of recent actions.",    ACCENT2),
        ("3. COMBO",  "Three actions form a combo (matches one of 5 moves).",      GOOD),
        ("4. RESOLVE","Every 5 sec, both queued moves resolve simultaneously.",    BAD),
    ]
    box_w = Inches(2.95)
    gap   = Inches(0.15)
    left  = Inches(0.5)
    top   = Inches(2.0)

    for i, (head, body, col) in enumerate(steps):
        x = left + (box_w + gap) * i
        _panel(s, x, top, box_w, Inches(2.6))
        _accent_bar(s, x + Inches(0.2), top + Inches(0.25), Inches(0.08), Inches(2.1),
                    color=col)
        _add_text(s, x + Inches(0.4), top + Inches(0.25), box_w - Inches(0.6),
                  Inches(0.6), head, size=20, bold=True, color=col)
        _add_text(s, x + Inches(0.4), top + Inches(0.95), box_w - Inches(0.6),
                  Inches(1.55), body, size=14, color=TEXT)

    _add_text(s, Inches(0.5), Inches(5.0), Inches(12.3), Inches(0.5),
              "HP starts at 100. Damage and heals apply on each resolution tick.",
              size=16, color=MUTED, align=PP_ALIGN.CENTER)

    # Resolution callout
    _panel(s, Inches(2.0), Inches(5.7), Inches(9.3), Inches(1.2),
           color=RGBColor(0x24, 0x2F, 0x52))
    _add_text(s, Inches(2.2), Inches(5.8), Inches(9.0), Inches(0.4),
              "Key invariant", size=12, bold=True, color=ACCENT2)
    _add_text(s, Inches(2.2), Inches(6.15), Inches(9.0), Inches(0.7),
              "Both players' queued moves are resolved at the SAME instant — "
              "no one goes first.", size=15, color=TEXT)


def slide_five_moves(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s)
    _slide_header(s, "The Five Moves",
                  "Two attacks · two defences · one heal")

    moves = [
        ("Power Strike",  "Attack — Weak", "Any 3 identical signs",
         "−25 HP to opponent", ACCENT),
        ("Combo Blast",   "Attack — Strong", "Move 1 → 2 → 3",
         "−45 HP to opponent", BAD),
        ("Shield Wall",   "Defence", "Move 3 → 4 → 5",
         "Blocks Power Strike (reflects 30%)\nLeaks 20% of Combo Blast", ACCENT2),
        ("Dodge Roll",    "Defence", "Move 1 → 3 → 5",
         "70% dodge vs Power Strike\n33% dodge vs Combo Blast", ACCENT2),
        ("Mend",          "Heal",    "X → Y → X   (X ≠ Y)",
         "+35 HP — cancelled if opponent attacks", GOOD),
    ]

    card_w = Inches(2.45)
    gap    = Inches(0.1)
    left   = Inches(0.4)
    top    = Inches(1.85)
    card_h = Inches(4.9)

    for i, (name, kind, combo, effect, col) in enumerate(moves):
        x = left + (card_w + gap) * i
        _panel(s, x, top, card_w, card_h)
        _accent_bar(s, x, top, card_w, Inches(0.12), color=col)
        _add_text(s, x + Inches(0.2), top + Inches(0.25), card_w - Inches(0.4),
                  Inches(0.55), name, size=18, bold=True, color=col)
        _add_text(s, x + Inches(0.2), top + Inches(0.85), card_w - Inches(0.4),
                  Inches(0.4), kind, size=12, color=MUTED)
        _add_text(s, x + Inches(0.2), top + Inches(1.35), card_w - Inches(0.4),
                  Inches(0.3), "Combo", size=11, bold=True, color=ACCENT)
        _add_text(s, x + Inches(0.2), top + Inches(1.65), card_w - Inches(0.4),
                  Inches(0.9), combo, size=14, color=TEXT)
        _add_text(s, x + Inches(0.2), top + Inches(2.7), card_w - Inches(0.4),
                  Inches(0.3), "Effect", size=11, bold=True, color=ACCENT)
        _add_text(s, x + Inches(0.2), top + Inches(3.0), card_w - Inches(0.4),
                  Inches(1.8), effect, size=14, color=TEXT)


def slide_resolution(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s)
    _slide_header(s, "Resolution Logic",
                  "How damage, defence, and healing interact each tick")

    # Step list on the left
    _panel(s, Inches(0.5), Inches(1.7), Inches(6.3), Inches(5.2))
    _add_text(s, Inches(0.8), Inches(1.9), Inches(5.8), Inches(0.4),
              "Each 5-second tick", size=18, bold=True, color=ACCENT)
    _add_bullets(s, Inches(0.8), Inches(2.4), Inches(5.8), Inches(4.5), [
        "Pop one move from each player's queue",
        "Roll attack / heal values (base ± 5 each round)",
        "If opponent attacked, opposing Mend is cancelled",
        "Shield Wall: full block vs Power Strike (30% reflect); 20% leak vs Combo Blast",
        "Dodge Roll: 70% dodge vs Power Strike; 33% vs Combo Blast",
        "Net HP changes apply to BOTH players at the same instant",
    ], size=15)

    # Key matchups on the right
    _panel(s, Inches(7.0), Inches(1.7), Inches(5.8), Inches(5.2))
    _add_text(s, Inches(7.3), Inches(1.9), Inches(5.3), Inches(0.4),
              "Key matchups", size=18, bold=True, color=ACCENT2)

    pairs = [
        ("Power Strike  vs  Power Strike", "Both −25 HP",          BAD),
        ("Combo Blast   vs  Combo Blast",  "Both −45 HP",          BAD),
        ("Power Strike  vs  Shield Wall",  "Attacker −7 HP (reflect)", ACCENT2),
        ("Combo Blast   vs  Shield Wall",  "Defender −9 HP (20% leak)", ACCENT2),
        ("Power Strike  vs  Dodge Roll",   "70% dodge OR full damage", ACCENT2),
        ("Mend          vs  (no attack)",  "+35 HP",               GOOD),
        ("Mend          vs  Attack",       "Heal cancelled — eat damage", BAD),
    ]
    y = 2.5
    for label, result, col in pairs:
        _add_text(s, Inches(7.3), Inches(y), Inches(3.4), Inches(0.35),
                  label, size=13, color=TEXT)
        _add_text(s, Inches(10.7), Inches(y), Inches(2.1), Inches(0.35),
                  result, size=13, bold=True, color=col)
        y += 0.42


def slide_controls(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s)
    _slide_header(s, "Controls — Keyboard Fallback",
                  "Always available; same five actions as the camera")

    _panel(s, Inches(2.5), Inches(2.0), Inches(8.3), Inches(4.5))
    _add_text(s, Inches(2.5), Inches(2.15), Inches(8.3), Inches(0.5),
              "P1 / P2 key bindings", size=20, bold=True, color=ACCENT,
              align=PP_ALIGN.CENTER)

    header = [("P1 key", ACCENT), ("P2 key", ACCENT2), ("Action", TEXT)]
    rows = [
        ("Q", "A", "Move 1"),
        ("W", "S", "Move 2"),
        ("E", "D", "Move 3"),
        ("R", "F", "Move 4"),
        ("T", "G", "Move 5"),
    ]
    cols = [Inches(3.4), Inches(5.8), Inches(8.2)]
    y = 3.0
    for i, (h, col) in enumerate(header):
        _add_text(s, cols[i], Inches(y), Inches(2.4), Inches(0.4),
                  h, size=16, bold=True, color=col, align=PP_ALIGN.CENTER)
    y = 3.5
    for p1, p2, act in rows:
        _add_text(s, cols[0], Inches(y), Inches(2.4), Inches(0.45),
                  p1, size=18, bold=True, color=TEXT, align=PP_ALIGN.CENTER)
        _add_text(s, cols[1], Inches(y), Inches(2.4), Inches(0.45),
                  p2, size=18, bold=True, color=TEXT, align=PP_ALIGN.CENTER)
        _add_text(s, cols[2], Inches(y), Inches(2.4), Inches(0.45),
                  act, size=16, color=TEXT, align=PP_ALIGN.CENTER)
        y += 0.5

    _add_text(s, Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.4),
              "Tip: practise on keyboard first to learn the combo timing, "
              "then plug in your model.",
              size=14, color=MUTED, align=PP_ALIGN.CENTER)


def slide_ai_overview(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s)
    _slide_header(s, "AI Hand-Sign Detection",
                  "How the camera turns gestures into game actions")

    _panel(s, Inches(0.5), Inches(1.7), Inches(12.3), Inches(2.4))
    flow = [
        ("Webcam frame", ACCENT),
        ("Preprocess",   ACCENT2),
        ("Your model",   GOOD),
        ("Top class + conf.", ACCENT2),
        ("Action → game", ACCENT),
    ]
    pad = Inches(0.4)
    arrow_w = Inches(0.5)
    avail = Inches(11.5) - arrow_w * (len(flow) - 1)
    box_w = avail / len(flow)
    x = Inches(0.9)
    y = Inches(2.4)
    for i, (label, col) in enumerate(flow):
        _panel(s, x, y, box_w, Inches(1.1), color=RGBColor(0x24, 0x2F, 0x52))
        _accent_bar(s, x, y, box_w, Inches(0.1), color=col)
        _add_text(s, x + Inches(0.1), y + Inches(0.32), box_w - Inches(0.2),
                  Inches(0.7), label, size=13, bold=True, color=TEXT,
                  align=PP_ALIGN.CENTER)
        if i < len(flow) - 1:
            ax = x + box_w + Inches(0.05)
            ay = y + Inches(0.45)
            arrow = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, ax, ay,
                                       arrow_w - Inches(0.1), Inches(0.2))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = ACCENT
            arrow.line.fill.background()
        x = x + box_w + arrow_w

    _panel(s, Inches(0.5), Inches(4.3), Inches(12.3), Inches(2.6))
    _add_text(s, Inches(0.8), Inches(4.45), Inches(11.6), Inches(0.4),
              "What the game reads from your team folder", size=18,
              bold=True, color=ACCENT)
    _add_bullets(s, Inches(0.8), Inches(4.95), Inches(11.6), Inches(2.0), [
        "team.env  — name, colour, MODEL_TYPE (landmark | cnn), per-class confidence thresholds",
        "models/  — trained weights (hand_sign_classifier.pkl + label_encoder.pkl  OR  hand_sign_cnn.pth)",
        "model_arch.py — CNN only, optional; lets your team ship a custom head architecture",
    ], size=15)


def slide_two_pipelines(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s)
    _slide_header(s, "Pick Your Pipeline",
                  "Two options — each team picks one in team.env")

    # Landmark card
    _panel(s, Inches(0.5), Inches(1.7), Inches(6.1), Inches(5.2))
    _accent_bar(s, Inches(0.5), Inches(1.7), Inches(6.1), Inches(0.12),
                color=ACCENT2)
    _add_text(s, Inches(0.8), Inches(1.95), Inches(5.6), Inches(0.5),
              "LANDMARK", size=24, bold=True, color=ACCENT2)
    _add_text(s, Inches(0.8), Inches(2.5), Inches(5.6), Inches(0.4),
              "MediaPipe keypoints → classical ML",
              size=14, color=MUTED)
    _add_bullets(s, Inches(0.8), Inches(3.0), Inches(5.6), Inches(3.8), [
        "MediaPipe extracts 21 (x,y,z) keypoints per hand",
        "126-dim feature vector → RandomForest or MLP",
        "Trains in seconds, even on CPU",
        "Works with a small dataset (~100 samples / class)",
        "Default & recommended for first-time teams",
    ], size=14, accent_color=ACCENT2)
    _add_text(s, Inches(0.8), Inches(6.2), Inches(5.6), Inches(0.6),
              "Best when: gestures are distinguishable by joint geometry.",
              size=12, color=MUTED)

    # CNN card
    _panel(s, Inches(6.8), Inches(1.7), Inches(6.0), Inches(5.2))
    _accent_bar(s, Inches(6.8), Inches(1.7), Inches(6.0), Inches(0.12),
                color=GOOD)
    _add_text(s, Inches(7.1), Inches(1.95), Inches(5.5), Inches(0.5),
              "CNN", size=24, bold=True, color=GOOD)
    _add_text(s, Inches(7.1), Inches(2.5), Inches(5.5), Inches(0.4),
              "Raw 224×224 frame → MobileNetV2 transfer learning",
              size=14, color=MUTED)
    _add_bullets(s, Inches(7.1), Inches(3.0), Inches(5.5), Inches(3.8), [
        "Fine-tunes a pretrained MobileNetV2 backbone",
        "Outputs softmax over the 5 move classes",
        "GPU recommended — use Google Colab for training",
        "Needs more samples (200+ per class is comfortable)",
        "Custom architectures supported via Teams/<N>/model_arch.py",
    ], size=14, accent_color=GOOD)
    _add_text(s, Inches(7.1), Inches(6.2), Inches(5.5), Inches(0.6),
              "Best when: gestures look similar by joints but differ visually "
              "(palm vs back, finger curl).",
              size=12, color=MUTED)


def slide_prep_overview(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s)
    _slide_header(s, "Team Preparation Roadmap",
                  "Five steps from kickoff to tournament-ready")

    steps = [
        ("1", "CONFIGURE", "Edit team.env — set NAME, COLOR, MODEL_TYPE",       ACCENT),
        ("2", "COLLECT",   "Run collect_data.exe — record samples for moves 1–5", ACCENT2),
        ("3", "TRAIN",     "Upload zip to Colab notebook — train + download weights", GOOD),
        ("4", "DEPLOY",    "Unzip the bundle into Teams/Team<N>/",              ACCENT),
        ("5", "TEST",      "Play practice matches; tune thresholds; retrain if needed", BAD),
    ]
    top = Inches(2.0)
    row_h = Inches(0.95)
    for i, (num, head, body, col) in enumerate(steps):
        y = top + row_h * i
        _panel(s, Inches(0.6), y, Inches(12.1), Inches(0.85))
        circle = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.85), y + Inches(0.13),
                                    Inches(0.6), Inches(0.6))
        circle.fill.solid()
        circle.fill.fore_color.rgb = col
        circle.line.fill.background()
        _add_text(s, Inches(0.85), y + Inches(0.18), Inches(0.6), Inches(0.5),
                  num, size=20, bold=True, color=BG, align=PP_ALIGN.CENTER)
        _add_text(s, Inches(1.65), y + Inches(0.15), Inches(2.5), Inches(0.5),
                  head, size=18, bold=True, color=col)
        _add_text(s, Inches(4.3), y + Inches(0.18), Inches(8.2), Inches(0.55),
                  body, size=15, color=TEXT)


def slide_step1_configure(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s)
    _slide_header(s, "Step 1 — Configure team.env",
                  "Each team has its own file at Teams/Team<N>/team.env")

    _panel(s, Inches(0.5), Inches(1.7), Inches(6.4), Inches(5.2))
    _add_text(s, Inches(0.8), Inches(1.85), Inches(5.8), Inches(0.5),
              "Example file", size=18, bold=True, color=ACCENT)
    code = (
        "NAME=ALPHA\n"
        "COLOR=255,120,60\n"
        "\n"
        "# landmark | cnn\n"
        "MODEL_TYPE=landmark\n"
        "\n"
        "# Confidence floor (0.0–1.0)\n"
        "THRESHOLD_MOVE1=0.6\n"
        "THRESHOLD_MOVE2=0.6\n"
        "THRESHOLD_MOVE3=0.6\n"
        "THRESHOLD_MOVE4=0.6\n"
        "THRESHOLD_MOVE5=0.6\n"
    )
    box = s.shapes.add_textbox(Inches(0.8), Inches(2.4), Inches(5.8), Inches(4.3))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.1)
    for i, line in enumerate(code.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run()
        r.text = line if line else " "
        r.font.name = "Consolas"
        r.font.size = Pt(14)
        col = MUTED if line.strip().startswith("#") else TEXT
        if "=" in line and not line.strip().startswith("#"):
            col = ACCENT2
        r.font.color.rgb = col

    _panel(s, Inches(7.1), Inches(1.7), Inches(5.7), Inches(5.2))
    _add_text(s, Inches(7.4), Inches(1.85), Inches(5.2), Inches(0.5),
              "What each field does", size=18, bold=True, color=ACCENT)
    _add_bullets(s, Inches(7.4), Inches(2.4), Inches(5.2), Inches(4.5), [
        "NAME — shown in lobby, scoreboard, end-of-match screen",
        "COLOR — RGB triple; tints HP bar, particles, UI accents",
        "MODEL_TYPE — landmark (default) or cnn; picks the inference path",
        "THRESHOLD_MOVE<N> — minimum confidence to register that class; raise to suppress false positives",
        "Save with UTF-8 encoding, no quotes around values",
    ], size=14)


def slide_step2_collect(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s)
    _slide_header(s, "Step 2 — Collect Training Data",
                  "Use the bundled collect_data.exe — no Python install needed")

    _panel(s, Inches(0.5), Inches(1.7), Inches(6.4), Inches(5.2))
    _add_text(s, Inches(0.8), Inches(1.85), Inches(5.8), Inches(0.5),
              "How to record", size=18, bold=True, color=ACCENT)
    _add_bullets(s, Inches(0.8), Inches(2.4), Inches(5.8), Inches(4.5), [
        "Double-click collect_data.exe (or run from terminal)",
        "Pick mode (cnn or landmark) and team number",
        "Stand in good lighting; both hands visible to the camera",
        "Press 1 / 2 / 3 / 4 / 5 to start recording for that move",
        "Press SPACE to stop, Q to quit and save",
        "Output lands in teams/Team<N>/ next to the exe",
    ], size=14)

    _panel(s, Inches(7.1), Inches(1.7), Inches(5.7), Inches(5.2))
    _add_text(s, Inches(7.4), Inches(1.85), Inches(5.2), Inches(0.5),
              "Best practices", size=18, bold=True, color=ACCENT2)
    _add_bullets(s, Inches(7.4), Inches(2.4), Inches(5.2), Inches(4.5), [
        "Aim for 200+ samples per class (more is better)",
        "Vary the angle, distance, and pose slightly each take",
        "Have different teammates record — model learns hand diversity",
        "Capture both palm-front and palm-back if your gestures use them (CNN)",
        "Mix lighting conditions to reduce brittleness",
        "Save the resulting teams/ folder — you'll zip it next",
    ], size=14, accent_color=ACCENT2)


def slide_step3_train(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s)
    _slide_header(s, "Step 3 — Train on Google Colab",
                  "code/model/sigil_strike_colab.ipynb walks through every step")

    _panel(s, Inches(0.5), Inches(1.7), Inches(12.3), Inches(5.2))

    items = [
        ("Setup",          "Installs deps, picks TEAM_NUM + MODEL_TYPE, checks for GPU"),
        ("Hyperparameters","Separate cells for CNN and landmark knobs — each one is explained"),
        ("Data",           "Upload your zipped teams/Team<N>/ folder (recommended path)"),
        ("Train",          "Fine-tunes MobileNetV2 (CNN) or fits RandomForest / MLP (landmark)"),
        ("Inference test", "Sanity-check predictions on uploaded images / Colab webcam"),
        ("Download",       "Bundles weights (+ model_arch.py for CNN) into a zip to your PC"),
    ]
    y = 1.95
    for label, body in items:
        _add_text(s, Inches(0.85), Inches(y), Inches(2.5), Inches(0.55),
                  label, size=17, bold=True, color=ACCENT)
        _add_text(s, Inches(3.4), Inches(y), Inches(9.2), Inches(0.55),
                  body, size=15, color=TEXT)
        y += 0.78

    _add_text(s, Inches(0.85), Inches(6.5), Inches(11.5), Inches(0.4),
              "Unzip the downloaded bundle straight into Teams/Team<N>/ — "
              "the game picks it up automatically.",
              size=14, color=ACCENT2)


def slide_step4_deploy(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s)
    _slide_header(s, "Step 4 — Deploy & Test",
                  "Drop the weights in, launch the game, play")

    _panel(s, Inches(0.5), Inches(1.7), Inches(12.3), Inches(2.4))
    _add_text(s, Inches(0.8), Inches(1.85), Inches(11.5), Inches(0.5),
              "Folder layout the game expects", size=18, bold=True, color=ACCENT)
    code = (
        "Teams/Team<N>/\n"
        "├── team.env\n"
        "├── model_arch.py            ← CNN only, optional\n"
        "└── models/\n"
        "    ├── hand_sign_cnn.pth                    (CNN)\n"
        "    └── hand_sign_classifier.pkl + label_encoder.pkl   (landmark)"
    )
    box = s.shapes.add_textbox(Inches(0.8), Inches(2.4), Inches(11.5), Inches(1.7))
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(code.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run()
        r.text = line if line else " "
        r.font.name = "Consolas"
        r.font.size = Pt(14)
        r.font.color.rgb = TEXT

    _panel(s, Inches(0.5), Inches(4.3), Inches(12.3), Inches(2.6))
    _add_text(s, Inches(0.8), Inches(4.45), Inches(11.5), Inches(0.5),
              "Test checklist", size=18, bold=True, color=ACCENT2)
    _add_bullets(s, Inches(0.8), Inches(4.95), Inches(11.5), Inches(2.0), [
        "Launch start_game.bat with two team IDs (e.g. start_game.bat 1 4) to skip the bracket",
        "Confirm the game prints your model loaded with the expected class names",
        "Throw each move 10× — note any low-confidence or misfire cases",
        "Tune THRESHOLD_MOVE<N> in team.env to suppress the most common false positives",
    ], size=14, accent_color=ACCENT2)


def slide_tips(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s)
    _slide_header(s, "Tips for a Strong Model",
                  "What separates a winning team's model from an average one")

    _panel(s, Inches(0.5), Inches(1.7), Inches(6.1), Inches(5.2))
    _accent_bar(s, Inches(0.5), Inches(1.7), Inches(6.1), Inches(0.12), color=GOOD)
    _add_text(s, Inches(0.8), Inches(1.95), Inches(5.6), Inches(0.5),
              "DO", size=22, bold=True, color=GOOD)
    _add_bullets(s, Inches(0.8), Inches(2.5), Inches(5.6), Inches(4.3), [
        "Pick 5 gestures that look visually distinct from each other",
        "Pick gestures you can hold steady for ~1 second",
        "Record with multiple teammates, multiple distances, multiple lighting",
        "Validate live before tournament day — not just on Colab",
        "Tune per-class thresholds — different gestures fail differently",
        "Re-collect if a class consistently misfires",
    ], size=14, accent_color=GOOD)

    _panel(s, Inches(6.8), Inches(1.7), Inches(6.0), Inches(5.2))
    _accent_bar(s, Inches(6.8), Inches(1.7), Inches(6.0), Inches(0.12), color=BAD)
    _add_text(s, Inches(7.1), Inches(1.95), Inches(5.5), Inches(0.5),
              "DON'T", size=22, bold=True, color=BAD)
    _add_bullets(s, Inches(7.1), Inches(2.5), Inches(5.5), Inches(4.3), [
        "Don't pick two gestures that differ only in finger micro-position",
        "Don't record all 200 samples from one angle / one person",
        "Don't skip the threshold tuning step — defaults are conservative",
        "Don't train only once — iterate after live testing",
        "Don't forget MODEL_TYPE in team.env must match your trained artefacts",
        "Don't rely on the camera alone — practise the keyboard fallback too",
    ], size=14, accent_color=BAD)


def slide_custom_arch(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s)
    _slide_header(s, "Optional: Custom CNN Architecture",
                  "Ship your own build_model() with the team")

    _panel(s, Inches(0.5), Inches(1.7), Inches(12.3), Inches(5.2))
    _add_text(s, Inches(0.8), Inches(1.85), Inches(11.5), Inches(0.5),
              "How it works", size=18, bold=True, color=ACCENT)
    _add_bullets(s, Inches(0.8), Inches(2.4), Inches(11.5), Inches(2.6), [
        "Place a model_arch.py with a build_model() function inside Teams/Team<N>/",
        "Both cnn/inference.py and cnn/train_model.py import and use it automatically",
        "The Colab notebook auto-generates one capturing your CNN_HIDDEN_DIM / dropouts",
        "If the file is absent, the default architecture is used — existing teams keep working",
    ], size=15)

    _add_text(s, Inches(0.8), Inches(5.25), Inches(11.5), Inches(0.5),
              "When to bother", size=18, bold=True, color=ACCENT2)
    _add_bullets(s, Inches(0.8), Inches(5.8), Inches(11.5), Inches(1.2), [
        "Default head overfits or underfits your dataset — try wider/narrower hidden dim",
        "You want a deeper classifier head or different dropout regime",
    ], size=15, accent_color=ACCENT2)


def slide_checklist(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s)
    _slide_header(s, "Tournament Day Checklist",
                  "Have these ready before your first match")

    _panel(s, Inches(0.5), Inches(1.7), Inches(12.3), Inches(5.2))
    items = [
        ("team.env",        "Final NAME, COLOR, MODEL_TYPE, thresholds set"),
        ("Trained weights", "models/ folder has the artefacts matching MODEL_TYPE"),
        ("Model arch",      "(CNN only) model_arch.py present if you customised the head"),
        ("Lighting",        "Test under the venue's lighting — bring a clip lamp if uncertain"),
        ("Backup",          "Both teammates know the keyboard fallback bindings (Q/W/E/R/T, A/S/D/F/G)"),
        ("Combos",          "Both players have memorised all 5 combo sequences"),
        ("Practice match",  "Played at least one full match against another team's setup"),
    ]
    y = 2.0
    for label, body in items:
        check = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.8), Inches(y + 0.05),
                                   Inches(0.35), Inches(0.35))
        check.fill.solid()
        check.fill.fore_color.rgb = GOOD
        check.line.fill.background()
        _add_text(s, Inches(0.8), Inches(y + 0.05), Inches(0.35), Inches(0.35),
                  "✓", size=14, bold=True, color=BG, align=PP_ALIGN.CENTER)
        _add_text(s, Inches(1.4), Inches(y), Inches(3.0), Inches(0.45),
                  label, size=16, bold=True, color=ACCENT)
        _add_text(s, Inches(4.5), Inches(y), Inches(8.0), Inches(0.45),
                  body, size=15, color=TEXT)
        y += 0.62


def slide_resources(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s)
    _slide_header(s, "Resources & Where to Look",
                  "Everything you need lives in the repo")

    _panel(s, Inches(0.5), Inches(1.7), Inches(12.3), Inches(5.2))

    rows = [
        ("README.md",                          "Full documentation, controls, resolution table"),
        ("start_game.bat",                     "Windows launcher — also accepts team IDs for test matches"),
        ("collect_data.exe",                   "Standalone data-collection tool (no Python required)"),
        ("code/model/sigil_strike_colab.ipynb","Colab notebook — open in colab.research.google.com"),
        ("Teams/Team<N>/team.env",             "Your team's config — name, colour, model type, thresholds"),
        ("code/configs/move_config.ini",       "Tune damage, heal, defence % — falls back to defaults if missing"),
        ("code/configs/time_config.ini",       "Tune resolve interval, deathmatch timing"),
    ]
    y = 1.95
    for label, body in rows:
        _add_text(s, Inches(0.8), Inches(y), Inches(4.4), Inches(0.5),
                  label, size=14, bold=True, color=ACCENT2, font="Consolas")
        _add_text(s, Inches(5.3), Inches(y), Inches(7.3), Inches(0.5),
                  body, size=14, color=TEXT)
        y += 0.65

    _add_text(s, Inches(0.5), Inches(6.95), Inches(12.3), Inches(0.4),
              "Good luck — may the best model win!",
              size=16, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)


# ── Driver ──────────────────────────────────────────────────────────────────────

def build(out_path: pathlib.Path) -> None:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    builders = [
        slide_title,
        slide_what_is_it,
        slide_tournament_format,
        slide_gameplay_loop,
        slide_five_moves,
        slide_resolution,
        slide_controls,
        slide_ai_overview,
        slide_two_pipelines,
        slide_prep_overview,
        slide_step1_configure,
        slide_step2_collect,
        slide_step3_train,
        slide_step4_deploy,
        slide_tips,
        slide_custom_arch,
        slide_checklist,
        slide_resources,
    ]
    total = len(builders)
    for i, fn in enumerate(builders, 1):
        fn(prs)
        _footer(prs.slides[i - 1], i, total)

    prs.save(out_path)
    print(f"[ok] wrote {out_path}  ({total} slides)")


if __name__ == "__main__":
    repo_root = pathlib.Path(__file__).resolve().parent.parent
    build(repo_root / "sigil_strike_team_briefing.pptx")
